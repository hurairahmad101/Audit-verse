"""Builds the AuditVerse.AI client pitch deck (PPTX) using python-pptx.

Dark navy theme styled to match the in-product UI. Each of the 17 audit
modules gets a slide that mirrors a real card/table/dashboard mockup
drawn natively with pptx shapes — same palette as the live app.
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Theme ----------
BG          = RGBColor(0x02, 0x06, 0x17)   # slate-950
PANEL       = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
PANEL_SOFT  = RGBColor(0x12, 0x1B, 0x32)
BORDER      = RGBColor(0x1F, 0x2A, 0x44)   # slate-700/60
BORDER_SOFT = RGBColor(0x18, 0x22, 0x39)
WHITE       = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_300   = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_400   = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500   = RGBColor(0x64, 0x74, 0x8B)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)
BLUE_SOFT   = RGBColor(0x1E, 0x3A, 0x8A)
EMERALD     = RGBColor(0x10, 0xB9, 0x81)
EMERALD_SOFT= RGBColor(0x06, 0x4E, 0x3B)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT  = RGBColor(0x78, 0x4F, 0x0A)
RED         = RGBColor(0xEF, 0x44, 0x44)
RED_SOFT    = RGBColor(0x7F, 0x1D, 0x1D)
PURPLE      = RGBColor(0xA8, 0x5C, 0xF6)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)

FONT_DISPLAY = "Calibri"   # PPTX-safe; bold/black weights for hero
FONT_BODY    = "Calibri"

# 16:9, 13.333" x 7.5" (standard widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Helpers ----------
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def text_box(slide, x, y, w, h, text, *, font=FONT_BODY, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, *, fill=PANEL, border=BORDER, border_w=0.75,
         corner=0.02, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = corner
    return sh


def pill(slide, x, y, text, *, fill=BLUE_SOFT, color=BLUE, w=Inches(1.2),
         h=Inches(0.28), size=9, bold=True):
    sh = rect(slide, x, y, w, h, fill=fill, border=fill, corner=0.5)
    tb = text_box(slide, x, y, w, h, text, size=size, bold=bold, color=color,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sh


def chip(slide, x, y, text, *, color=SLATE_400, size=9):
    """Inline label without background, e.g. metadata."""
    return text_box(slide, x, y, Inches(2.5), Inches(0.25), text, size=size,
                    color=color)


def divider(slide, x, y, w, color=BORDER, weight=0.5):
    sh = slide.shapes.add_connector(1, x, y, x + w, y)
    sh.line.color.rgb = color
    sh.line.width = Pt(weight)
    return sh


def header_band(slide, eyebrow, title, subtitle):
    """Standard module-slide header."""
    text_box(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
             eyebrow.upper(), font=FONT_DISPLAY, size=10, bold=True,
             color=BLUE)
    text_box(slide, Inches(0.6), Inches(0.75), Inches(11), Inches(0.7),
             title, font=FONT_DISPLAY, size=32, bold=True, color=WHITE)
    text_box(slide, Inches(0.6), Inches(1.4), Inches(11), Inches(0.4),
             subtitle, size=13, color=SLATE_400)
    divider(slide, Inches(0.6), Inches(1.95), Inches(12.15))


def footer(slide, page_num, total):
    text_box(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "AuditVerse.AI  ·  Audit Management Platform", size=9,
             color=SLATE_500)
    text_box(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page_num} / {total}", size=9, color=SLATE_500,
             align=PP_ALIGN.RIGHT)


# ---------- Reusable mockup panels ----------
def stat_card(slide, x, y, w, h, label, value, *, accent=BLUE,
              sub=None):
    rect(slide, x, y, w, h, fill=PANEL, border=BORDER)
    # accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18),
                                  y + Inches(0.22), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    text_box(slide, x + Inches(0.42), y + Inches(0.18), w - Inches(0.5),
             Inches(0.3), label, size=9, color=SLATE_400, bold=True)
    text_box(slide, x + Inches(0.18), y + Inches(0.55), w - Inches(0.3),
             Inches(0.6), value, size=24, bold=True, color=WHITE)
    if sub:
        text_box(slide, x + Inches(0.18), y + h - Inches(0.4),
                 w - Inches(0.3), Inches(0.3), sub, size=9, color=SLATE_500)


def severity_pill(slide, x, y, severity):
    palette = {
        "critical": (RED_SOFT, RED, "Critical"),
        "high":     (AMBER_SOFT, AMBER, "High"),
        "medium":   (RGBColor(0x78, 0x59, 0x0A), AMBER, "Medium"),
        "low":      (EMERALD_SOFT, EMERALD, "Low"),
        "open":     (RED_SOFT, RED, "Open"),
        "in_progress": (AMBER_SOFT, AMBER, "In Progress"),
        "remediated": (EMERALD_SOFT, EMERALD, "Remediated"),
        "closed":   (RGBColor(0x33, 0x41, 0x55), SLATE_400, "Closed"),
        "draft":    (RGBColor(0x33, 0x41, 0x55), SLATE_400, "Draft"),
        "approved": (EMERALD_SOFT, EMERALD, "Approved"),
        "review":   (RGBColor(0x1E, 0x40, 0xAF), BLUE, "In Review"),
        "signed":   (RGBColor(0x4C, 0x1D, 0x95), PURPLE, "Signed Off"),
    }
    bg, fg, label = palette.get(severity.lower(), (PANEL, SLATE_400, severity))
    pill(slide, x, y, label, fill=bg, color=fg, w=Inches(0.95), h=Inches(0.26),
         size=8)


# ---------- Slide builders ----------
def make_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG)
    # Large gradient panel on right
    panel = rect(s, Inches(7.2), Inches(0.5), Inches(5.6), Inches(6.5),
                 fill=PANEL_SOFT, border=BORDER)
    # Mock product chrome inside
    rect(s, Inches(7.5), Inches(0.85), Inches(5.0), Inches(0.45),
         fill=BG, border=BORDER_SOFT, corner=0.15)
    text_box(s, Inches(7.7), Inches(0.93), Inches(2), Inches(0.3),
             "AuditVerse.AI", size=11, bold=True, color=WHITE)
    text_box(s, Inches(11.4), Inches(0.93), Inches(1), Inches(0.3),
             "Live", size=9, color=EMERALD)
    # Mini cards grid 2x3
    labels = [("Open Findings", "47", RED), ("Engagements", "23", BLUE),
              ("Closure Rate", "82%", EMERALD), ("Overdue", "8", AMBER),
              ("Coverage", "91%", PURPLE), ("Cycle Time", "18d", CYAN)]
    for i, (lab, val, col) in enumerate(labels):
        cx = Inches(7.5 + (i % 2) * 2.55)
        cy = Inches(1.5 + (i // 2) * 1.7)
        stat_card(s, cx, cy, Inches(2.4), Inches(1.5), lab, val, accent=col)

    # Left side hero
    text_box(s, Inches(0.7), Inches(1.0), Inches(2), Inches(0.35),
             "AUDITVERSE.AI", size=10, bold=True, color=BLUE)
    text_box(s, Inches(0.7), Inches(1.4), Inches(6.5), Inches(2.4),
             "The audit\nplatform that\nactually ships\nfindings.",
             font=FONT_DISPLAY, size=54, bold=True, color=WHITE,
             line_spacing=0.95)
    text_box(s, Inches(0.7), Inches(4.8), Inches(6.3), Inches(1.4),
             ("One workspace for the audit lifecycle — universe to "
              "engagement, workpaper to finding, escalation to closure."),
             size=16, color=SLATE_300, line_spacing=1.35)
    text_box(s, Inches(0.7), Inches(6.6), Inches(6), Inches(0.4),
             "Client Briefing  ·  May 2026", size=11, color=SLATE_400)


def make_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    text_box(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
             "THE PROBLEM", size=10, bold=True, color=AMBER)
    text_box(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9),
             "Internal audit runs on email, spreadsheets, and stale memory.",
             font=FONT_DISPLAY, size=30, bold=True, color=WHITE)
    text_box(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
             ("Three groups feel the pain — and they don't see the same "
              "version of the truth."),
             size=14, color=SLATE_400)

    cols = [
        ("EXECUTIVES", BLUE,
         "No real-time view of audit coverage, open risks or remediation status.",
         ["Last quarter's deck is the only artifact",
          "Audit committee learns of issues 60+ days late",
          "Assurance gaps invisible until something breaks"]),
        ("AUDIT TEAM", AMBER,
         "Spend 60% of cycle time copy-pasting between Word, Excel and email.",
         ["Workpaper review threads buried in inbox",
          "No standard CCCE structure across findings",
          "Skill-to-engagement matching done in heads"]),
        ("BUSINESS / IT", EMERALD,
         "Repeated requests for the same evidence, no visibility into status.",
         ["Same control re-audited from 3 angles",
          "Action owners chase audit for due dates",
          "No self-service to attest or push back"]),
    ]
    for i, (head, col, lead, bullets) in enumerate(cols):
        x = Inches(0.6 + i * 4.18)
        rect(s, x, Inches(2.4), Inches(4.0), Inches(4.4),
             fill=PANEL, border=BORDER)
        text_box(s, x + Inches(0.3), Inches(2.6), Inches(3.6), Inches(0.3),
                 head, size=10, bold=True, color=col)
        text_box(s, x + Inches(0.3), Inches(2.95), Inches(3.6), Inches(1.2),
                 lead, size=14, bold=True, color=WHITE, line_spacing=1.25)
        for j, b in enumerate(bullets):
            ty = Inches(4.3 + j * 0.7)
            # bullet dot
            d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x + Inches(0.3), ty + Inches(0.1),
                                   Inches(0.08), Inches(0.08))
            d.fill.solid(); d.fill.fore_color.rgb = col; d.line.fill.background()
            text_box(s, x + Inches(0.5), ty, Inches(3.4), Inches(0.6),
                     b, size=11, color=SLATE_300, line_spacing=1.35)


def make_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    text_box(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
             "THE SOLUTION", size=10, bold=True, color=EMERALD)
    text_box(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9),
             "Seventeen modules. One workspace. The full audit lifecycle.",
             font=FONT_DISPLAY, size=28, bold=True, color=WHITE)
    text_box(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             "Each module is on the next slide. Here is the map.",
             size=13, color=SLATE_400)

    modules = [
        ("Universe", BLUE), ("Plans", BLUE), ("Engagements", BLUE),
        ("Workpapers", AMBER), ("Findings", RED), ("Issues", RED),
        ("Surveys", PURPLE), ("Documents", CYAN), ("Analytics", EMERALD),
        ("Portal", PURPLE), ("Charter", CYAN), ("CCM", EMERALD),
        ("Reporting", BLUE), ("QAIP", AMBER), ("Test Scripts", AMBER),
        ("Skill Matrix", PURPLE), ("Capacity", PURPLE),
    ]
    cols = 6
    cell_w = Inches(2.05); cell_h = Inches(0.95)
    gap = Inches(0.1)
    grid_x0 = Inches(0.6); grid_y0 = Inches(2.4)
    for i, (name, col) in enumerate(modules):
        r, c = divmod(i, cols)
        x = grid_x0 + c * (cell_w + gap)
        y = grid_y0 + r * (cell_h + gap)
        rect(s, x, y, cell_w, cell_h, fill=PANEL, border=BORDER)
        # accent stripe top
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y, cell_w, Inches(0.06))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        text_box(s, x + Inches(0.15), y + Inches(0.18), cell_w - Inches(0.3),
                 Inches(0.3), f"{i+1:02d}", size=9, color=SLATE_500, bold=True)
        text_box(s, x + Inches(0.15), y + Inches(0.42), cell_w - Inches(0.3),
                 Inches(0.5), name, size=14, bold=True, color=WHITE)


# ---------- Module slide template ----------
def module_slide(prs, eyebrow, title, problem, solution, mockup_fn,
                 features, page_num, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header_band(s, eyebrow, title, problem["lead"])

    # LEFT column: problem + solution narrative
    text_box(s, Inches(0.6), Inches(2.15), Inches(0.8), Inches(0.3),
             "TODAY", size=9, bold=True, color=AMBER)
    text_box(s, Inches(0.6), Inches(2.45), Inches(5.2), Inches(1.4),
             problem["body"], size=12, color=SLATE_300, line_spacing=1.4)

    text_box(s, Inches(0.6), Inches(3.85), Inches(2), Inches(0.3),
             "WITH AUDITVERSE.AI", size=9, bold=True, color=EMERALD)
    text_box(s, Inches(0.6), Inches(4.15), Inches(5.2), Inches(1.5),
             solution, size=12, color=SLATE_300, line_spacing=1.4)

    # Feature chips
    fy = Inches(5.7)
    for i, f in enumerate(features[:4]):
        ty = fy + i * Inches(0.32)
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), ty + Inches(0.07),
                                 Inches(0.08), Inches(0.16))
        sh.fill.solid(); sh.fill.fore_color.rgb = BLUE
        sh.line.fill.background()
        text_box(s, Inches(0.85), ty, Inches(5.0), Inches(0.3),
                 f, size=10, color=SLATE_300)

    # RIGHT column: mockup
    rect(s, Inches(6.2), Inches(2.15), Inches(6.55), Inches(4.7),
         fill=PANEL, border=BORDER)
    # Window chrome dots
    for i, c in enumerate([RGBColor(0xEF, 0x44, 0x44),
                            RGBColor(0xF5, 0x9E, 0x0B),
                            RGBColor(0x10, 0xB9, 0x81)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(6.4 + i * 0.22), Inches(2.32),
                               Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
    text_box(s, Inches(7.3), Inches(2.27), Inches(5), Inches(0.3),
             f"auditverse.ai  ·  /audit/{eyebrow.lower().replace(' ', '-')}",
             size=8, color=SLATE_500)
    divider(s, Inches(6.2), Inches(2.62), Inches(6.55), color=BORDER_SOFT)

    # Mockup body area — prefer a real screenshot of the live app if one is
    # available on disk, otherwise fall back to the procedural mockup. Real
    # screenshots are captured by scripts/capture_screenshots.py at
    # 1920x1080 (16:9), so we letterbox the image inside the panel slot
    # (6.15w x 3.95h) to preserve aspect ratio without cropping content.
    import os as _os
    slug = eyebrow.lower().replace(" ", "-")
    shot = _os.path.join("assets", "pitch", "screenshots", f"{slug}.png")
    if _os.path.exists(shot):
        slot_x, slot_y = Inches(6.4), Inches(2.78)
        slot_w, slot_h = Inches(6.15), Inches(3.95)
        img_aspect = 1920 / 1080  # 16:9
        slot_aspect = slot_w / slot_h
        if img_aspect >= slot_aspect:
            # Width-bound (wider than slot) — fit by width, center vertically
            pic_w = slot_w
            pic_h = int(slot_w / img_aspect)
        else:
            # Height-bound (taller than slot) — fit by height, center horizontally
            pic_h = slot_h
            pic_w = int(slot_h * img_aspect)
        pic_x = slot_x + (slot_w - pic_w) / 2
        pic_y = slot_y + (slot_h - pic_h) / 2
        s.shapes.add_picture(shot, pic_x, pic_y, width=pic_w, height=pic_h)
    else:
        mockup_fn(s, Inches(6.4), Inches(2.78), Inches(6.15), Inches(3.95))

    footer(s, page_num, total)


# ---------- Mockup variants per module ----------
def m_universe(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Audit Universe — Risk Heatmap",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "32 entities · 7 business units", size=9, color=SLATE_500)
    # 5x4 heat grid
    cols = 8; rows = 4
    cell_w = (w - Inches(0.1)) / cols
    cell_h = Inches(0.55)
    palette = [EMERALD, EMERALD, AMBER, AMBER, RED, RED, BLUE,
               EMERALD, AMBER, RED, AMBER, EMERALD, BLUE, AMBER,
               RED, RED, AMBER, EMERALD, EMERALD, BLUE, AMBER, RED,
               AMBER, EMERALD, BLUE, EMERALD, AMBER, RED, EMERALD,
               EMERALD, AMBER, BLUE]
    labels = ["Treasury", "Lending", "Payments", "Cards", "AML",
              "Cyber", "Cloud", "HR", "Procure", "Treasury", "Lending",
              "Payments", "Cards", "AML", "Cyber", "Cloud",
              "Treasury", "Lending", "Payments", "Cards", "AML",
              "Cyber", "Cloud", "HR", "Treasury", "Lending", "Payments",
              "Cards", "AML", "Cyber", "Cloud", "HR"]
    for i in range(rows * cols):
        r, c = divmod(i, cols)
        cx = x + Emu(c * int(cell_w)) + Inches(0.05)
        cy = y + Inches(0.7) + r * (cell_h + Inches(0.06))
        col = palette[i % len(palette)]
        rect(s, cx, cy, cell_w - Inches(0.06), cell_h, fill=col, border=col,
             corner=0.18)
        text_box(s, cx, cy, cell_w - Inches(0.06), cell_h, labels[i % len(labels)],
                 size=8, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # legend
    ly = y + Inches(3.4)
    legend = [("Low", EMERALD), ("Medium", AMBER), ("High", RED),
              ("No coverage", BLUE)]
    for i, (lab, col) in enumerate(legend):
        lx = x + Inches(i * 1.45)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, ly + Inches(0.08),
                               Inches(0.14), Inches(0.14))
        d.fill.solid(); d.fill.fore_color.rgb = col; d.line.fill.background()
        text_box(s, lx + Inches(0.22), ly, Inches(1.3), Inches(0.3),
                 lab, size=9, color=SLATE_400)


def m_plans(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Annual Audit Plan — FY2026",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "23 engagements · 4 quarters · 14,200 budgeted hours",
             size=9, color=SLATE_500)
    # quarter columns with stacked engagements
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    items = [
        [("Cyber Resilience", BLUE), ("Treasury IT", EMERALD), ("AML KYC", AMBER)],
        [("Cards Fraud", RED), ("Cloud Migration", BLUE), ("HR Payroll", EMERALD)],
        [("Vendor Risk", AMBER), ("Branch Ops", BLUE), ("Data Privacy", PURPLE)],
        [("Year-end ITGC", BLUE), ("AML Refresh", AMBER), ("Cyber Followup", RED)],
    ]
    col_w = (w - Inches(0.3)) / 4
    for i, q in enumerate(quarters):
        cx = x + Inches(0.05) + i * col_w
        rect(s, cx, y + Inches(0.7), col_w - Inches(0.1), Inches(2.9),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, cx + Inches(0.1), y + Inches(0.78), col_w, Inches(0.3),
                 q, size=11, bold=True, color=BLUE)
        for j, (name, col) in enumerate(items[i]):
            ey = y + Inches(1.15 + j * 0.7)
            rect(s, cx + Inches(0.1), ey, col_w - Inches(0.3), Inches(0.6),
                 fill=PANEL, border=BORDER_SOFT)
            stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         cx + Inches(0.1), ey,
                                         Inches(0.05), Inches(0.6))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = col
            stripe.line.fill.background()
            text_box(s, cx + Inches(0.22), ey + Inches(0.08),
                     col_w - Inches(0.4), Inches(0.25),
                     name, size=9, bold=True, color=WHITE)
            text_box(s, cx + Inches(0.22), ey + Inches(0.32),
                     col_w - Inches(0.4), Inches(0.25),
                     f"{(j+1)*120}h budget", size=8, color=SLATE_500)


def _engagement_row(s, x, y, w, h, name, status_key, hours_used, hours_budget,
                    owner):
    rect(s, x, y, w, h, fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, x + Inches(0.15), y + Inches(0.08), w - Inches(2),
             Inches(0.3), name, size=10, bold=True, color=WHITE)
    chip(s, x + Inches(0.15), y + Inches(0.35), f"Owner: {owner}",
         color=SLATE_500, size=8)
    severity_pill(s, x + w - Inches(1.15), y + Inches(0.1), status_key)
    # progress bar
    pct = min(1.0, hours_used / max(1, hours_budget))
    bar_w = Inches(2.4)
    bx = x + w - Inches(1.15) - bar_w - Inches(0.2)
    rect(s, bx, y + Inches(0.45), bar_w, Inches(0.18),
         fill=BG, border=BORDER_SOFT, corner=0.5)
    fill_w = Emu(int(bar_w * pct))
    if fill_w > Emu(0):
        col = EMERALD if pct < 0.85 else (AMBER if pct < 1.0 else RED)
        rect(s, bx, y + Inches(0.45), fill_w, Inches(0.18),
             fill=col, border=col, corner=0.5)
    text_box(s, bx, y + Inches(0.05), bar_w, Inches(0.3),
             f"{hours_used} / {hours_budget} hrs",
             size=8, color=SLATE_400, align=PP_ALIGN.RIGHT)


def m_engagements(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Active Engagements",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "8 in-flight  ·  3 in planning  ·  2 in reporting",
             size=9, color=SLATE_500)
    rows = [
        ("Cyber Resilience Review", "in_progress", 412, 480, "P. Singh"),
        ("Treasury IT General Controls", "in_progress", 180, 320, "L. Chen"),
        ("AML KYC Walkthrough", "review", 295, 280, "R. Mendoza"),
        ("Cards Fraud Analytics", "draft", 0, 240, "A. Khan"),
        ("Vendor Risk Refresh", "remediated", 195, 200, "T. Olsson"),
    ]
    for i, row in enumerate(rows):
        ry = y + Inches(0.7 + i * 0.72)
        _engagement_row(s, x, ry, w - Inches(0.2), Inches(0.65), *row)


def m_workpapers(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Workpaper Sign-Off Workflow",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Cyber Resilience Review · 47 workpapers", size=9, color=SLATE_500)
    rows = [
        ("WP-001", "Walkthrough — Identity & Access", "approved", "P. Singh", "L. Chen"),
        ("WP-002", "Test of Privileged Access Reviews", "review", "P. Singh", "L. Chen"),
        ("WP-003", "Analytical — Failed Login Trend", "signed", "A. Khan", "L. Chen"),
        ("WP-004", "Interview — IR Lead", "draft", "P. Singh", "—"),
        ("WP-005", "Test — MFA Coverage", "approved", "A. Khan", "L. Chen"),
    ]
    # header row
    rect(s, x, y + Inches(0.7), w - Inches(0.2), Inches(0.35),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    headers = [("Ref", 0.15, 0.6), ("Title", 0.85, 2.6),
               ("Status", 3.55, 1.0), ("Preparer", 4.65, 0.9),
               ("Reviewer", 5.6, 0.9)]
    for h_lab, hx, hw in headers:
        text_box(s, x + Inches(hx), y + Inches(0.78),
                 Inches(hw), Inches(0.25),
                 h_lab.upper(), size=8, bold=True, color=SLATE_500)
    for i, (ref, title, status, prep, rev) in enumerate(rows):
        ry = y + Inches(1.1 + i * 0.5)
        if i % 2 == 0:
            rect(s, x, ry, w - Inches(0.2), Inches(0.45),
                 fill=PANEL_SOFT, border=PANEL_SOFT)
        text_box(s, x + Inches(0.15), ry + Inches(0.12), Inches(0.6),
                 Inches(0.3), ref, size=9, color=SLATE_400, bold=True)
        text_box(s, x + Inches(0.85), ry + Inches(0.12), Inches(2.6),
                 Inches(0.3), title, size=9, color=WHITE)
        severity_pill(s, x + Inches(3.55), ry + Inches(0.1), status)
        text_box(s, x + Inches(4.65), ry + Inches(0.12), Inches(0.9),
                 Inches(0.3), prep, size=9, color=SLATE_400)
        text_box(s, x + Inches(5.6), ry + Inches(0.12), Inches(0.9),
                 Inches(0.3), rev, size=9, color=SLATE_400)


def m_findings(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Findings — CCCE Structured",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "F-2026-014 · Cyber Resilience Review", size=9, color=SLATE_500)
    severity_pill(s, x + w - Inches(1.5), y + Inches(0.05), "critical")

    # CCCE 2x2
    cells = [
        ("CONDITION", "Privileged accounts not reviewed in 14 of 22 sampled apps.",
         RED),
        ("CRITERIA", "Internal Identity Policy IAM-007 requires quarterly review.",
         BLUE),
        ("CAUSE", "Quarterly review job disabled when SSO migrated in Feb.",
         AMBER),
        ("EFFECT", "21,400 dormant privileged sessions could be re-activated.",
         PURPLE),
    ]
    cw = (w - Inches(0.3)) / 2
    ch = Inches(1.35)
    for i, (lab, body, col) in enumerate(cells):
        r, c = divmod(i, 2)
        cx = x + c * (cw + Inches(0.1))
        cy = y + Inches(0.85) + r * (ch + Inches(0.1))
        rect(s, cx, cy, cw, ch, fill=PANEL_SOFT, border=BORDER_SOFT)
        # left accent stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cx, cy, Inches(0.06), ch)
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        text_box(s, cx + Inches(0.18), cy + Inches(0.1), cw - Inches(0.3),
                 Inches(0.3), lab, size=9, bold=True, color=col)
        text_box(s, cx + Inches(0.18), cy + Inches(0.4), cw - Inches(0.3),
                 ch - Inches(0.5), body, size=11, color=SLATE_300,
                 line_spacing=1.35)
    # AI suggestion strip
    sy = y + Inches(3.65)
    rect(s, x, sy, w - Inches(0.2), Inches(0.4),
         fill=RGBColor(0x1E, 0x1B, 0x4B), border=RGBColor(0x4C, 0x1D, 0x95))
    text_box(s, x + Inches(0.15), sy + Inches(0.08), Inches(2),
             Inches(0.25), "AI SUGGESTED",
             size=8, bold=True, color=PURPLE)
    text_box(s, x + Inches(1.4), sy + Inches(0.08), w - Inches(1.6),
             Inches(0.25),
             "Severity → Critical · Recommend re-enable review job + 30-day attestation",
             size=9, color=SLATE_300)


def m_issues(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Issue Tracking & Escalation",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Aging buckets · Active escalations", size=9, color=SLATE_500)
    # Aging buckets
    buckets = [("0–30d", 12, EMERALD), ("30–60d", 7, AMBER),
               ("60–90d", 4, RED), ("90d+", 2, RED)]
    bw = (w - Inches(0.4)) / 4
    for i, (lab, val, col) in enumerate(buckets):
        bx = x + i * (bw + Inches(0.1))
        rect(s, bx, y + Inches(0.7), bw, Inches(1.0),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, bx, y + Inches(0.78), bw, Inches(0.25),
                 lab, size=9, color=SLATE_500, align=PP_ALIGN.CENTER, bold=True)
        text_box(s, bx, y + Inches(1.05), bw, Inches(0.55), str(val),
                 size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Escalations list
    text_box(s, x, y + Inches(1.85), w, Inches(0.3),
             "ACTIVE ESCALATIONS", size=9, bold=True, color=SLATE_500)
    rows = [
        ("L3 — Executive", "MFA gaps in Cards platform", 92, RED),
        ("L2 — Manager", "Vendor SOC2 report missing", 47, AMBER),
        ("L2 — Manager", "Privileged review job disabled", 38, AMBER),
        ("L1 — First reminder", "Branch reconciliation backlog", 12, BLUE),
    ]
    for i, (lvl, ttl, days, col) in enumerate(rows):
        ry = y + Inches(2.2 + i * 0.42)
        rect(s, x, ry, w - Inches(0.2), Inches(0.36),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        # arrow up
        ar = s.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 x + Inches(0.15), ry + Inches(0.08),
                                 Inches(0.18), Inches(0.2))
        ar.fill.solid(); ar.fill.fore_color.rgb = col; ar.line.fill.background()
        text_box(s, x + Inches(0.4), ry + Inches(0.08), Inches(1.6),
                 Inches(0.25), lvl, size=9, bold=True, color=col)
        text_box(s, x + Inches(2.1), ry + Inches(0.08), Inches(3.2),
                 Inches(0.25), ttl, size=9, color=WHITE)
        text_box(s, x + Inches(5.4), ry + Inches(0.08), Inches(0.7),
                 Inches(0.25), f"{days}d", size=9, color=col, bold=True,
                 align=PP_ALIGN.RIGHT)


def m_surveys(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Surveys & Questionnaires",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Pre-audit · Post-audit · Control self-assessment",
             size=9, color=SLATE_500)
    cards = [
        ("Pre-Audit — IT Controls", "completed", 18, 18, BLUE),
        ("CSA — AML Operations", "in_progress", 9, 14, AMBER),
        ("Post-Audit — Treasury", "draft", 0, 8, SLATE_500),
    ]
    cw = (w - Inches(0.3)) / 3
    for i, (ttl, st, done, total, col) in enumerate(cards):
        cx = x + i * (cw + Inches(0.1))
        rect(s, cx, y + Inches(0.7), cw, Inches(2.3),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        severity_pill(s, cx + Inches(0.15), y + Inches(0.78), st)
        text_box(s, cx + Inches(0.15), y + Inches(1.15), cw - Inches(0.3),
                 Inches(0.6), ttl, size=11, bold=True, color=WHITE,
                 line_spacing=1.2)
        # response donut-ish: just %
        pct = int(100 * done / max(1, total))
        text_box(s, cx + Inches(0.15), y + Inches(1.85), cw - Inches(0.3),
                 Inches(0.55), f"{pct}%", size=24, bold=True, color=col)
        text_box(s, cx + Inches(0.15), y + Inches(2.45), cw - Inches(0.3),
                 Inches(0.3), f"{done} of {total} responses",
                 size=9, color=SLATE_500)
    # AI generate strip
    sy = y + Inches(3.15)
    rect(s, x, sy, w - Inches(0.2), Inches(0.4),
         fill=RGBColor(0x1E, 0x1B, 0x4B), border=RGBColor(0x4C, 0x1D, 0x95))
    text_box(s, x + Inches(0.15), sy + Inches(0.08), Inches(2),
             Inches(0.25), "AI GENERATED",
             size=8, bold=True, color=PURPLE)
    text_box(s, x + Inches(1.4), sy + Inches(0.08), w - Inches(1.6),
             Inches(0.25),
             "12 questions auto-drafted from engagement scope in 3 seconds",
             size=9, color=SLATE_300)


def m_documents(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Document Repository",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "412 documents · 1.4 GB · confidentiality flags",
             size=9, color=SLATE_500)
    stats = [("Total", "412", BLUE), ("Confidential", "47", RED),
             ("Storage", "1.4 GB", PURPLE), ("Types", "6", EMERALD)]
    sw = (w - Inches(0.3)) / 4
    for i, (lab, val, col) in enumerate(stats):
        sx = x + i * (sw + Inches(0.1))
        stat_card(s, sx, y + Inches(0.7), sw, Inches(1.05), lab, val,
                  accent=col)
    # File list
    files = [
        ("ITGC Walkthrough Notes.docx", "Working Paper", "Cyber Resilience", True),
        ("Sample Population FY26.xlsx", "Evidence", "AML KYC", False),
        ("IR Plan v3.pdf", "Evidence", "Cyber Resilience", True),
        ("Charter 2026 Approved.pdf", "Charter", "—", False),
    ]
    text_box(s, x, y + Inches(1.95), w, Inches(0.3),
             "RECENTLY ADDED", size=9, bold=True, color=SLATE_500)
    for i, (name, typ, eng, conf) in enumerate(files):
        ry = y + Inches(2.25 + i * 0.4)
        rect(s, x, ry, w - Inches(0.2), Inches(0.34),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, x + Inches(0.15), ry + Inches(0.08), Inches(3.0),
                 Inches(0.25), name, size=9, bold=True, color=WHITE)
        text_box(s, x + Inches(3.2), ry + Inches(0.08), Inches(1.4),
                 Inches(0.25), typ, size=9, color=BLUE)
        text_box(s, x + Inches(4.6), ry + Inches(0.08), Inches(1.6),
                 Inches(0.25), eng, size=9, color=SLATE_400)
        if conf:
            pill(s, x + Inches(5.2), ry + Inches(0.05),
                 "Confidential", fill=RED_SOFT, color=RED,
                 w=Inches(0.85), h=Inches(0.24), size=7)


def m_analytics(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Audit Analytics — Coverage & Velocity",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Rolling 12-month view", size=9, color=SLATE_500)
    # Bar chart mock — findings closed by quarter
    chart_x = x; chart_y = y + Inches(0.7)
    chart_w = w - Inches(0.2); chart_h = Inches(2.0)
    rect(s, chart_x, chart_y, chart_w, chart_h, fill=PANEL_SOFT,
         border=BORDER_SOFT)
    bars = [("Q1 25", 24, 18), ("Q2 25", 31, 27), ("Q3 25", 28, 26),
            ("Q4 25", 35, 33), ("Q1 26", 42, 38)]
    bx0 = chart_x + Inches(0.4); by_top = chart_y + Inches(0.25)
    avail_w = chart_w - Inches(0.6); avail_h = chart_h - Inches(0.65)
    bw = avail_w / len(bars) / 2.5
    for i, (lab, opened, closed) in enumerate(bars):
        slot = avail_w / len(bars)
        gx = bx0 + i * slot
        h_open = Emu(int(avail_h * opened / 50))
        h_clos = Emu(int(avail_h * closed / 50))
        rect(s, gx, by_top + avail_h - h_open, bw, h_open,
             fill=AMBER, border=AMBER, corner=0.2)
        rect(s, gx + bw + Inches(0.05), by_top + avail_h - h_clos, bw, h_clos,
             fill=EMERALD, border=EMERALD, corner=0.2)
        text_box(s, gx, by_top + avail_h + Inches(0.05), bw * 2.2,
                 Inches(0.25), lab, size=8, color=SLATE_400,
                 align=PP_ALIGN.CENTER)
    # legend
    text_box(s, chart_x + Inches(0.4), chart_y + Inches(0.05), Inches(2),
             Inches(0.25), "■ Opened", size=9, color=AMBER, bold=True)
    text_box(s, chart_x + Inches(1.6), chart_y + Inches(0.05), Inches(2),
             Inches(0.25), "■ Closed", size=9, color=EMERALD, bold=True)
    # KPI row
    kpis = [("Coverage", "91%", EMERALD), ("Mean cycle", "18d", BLUE),
            ("Closure rate", "82%", EMERALD), ("Findings/eng", "2.7", AMBER)]
    kw = (w - Inches(0.4)) / 4
    for i, (lab, val, col) in enumerate(kpis):
        kx = x + i * (kw + Inches(0.1))
        ky = y + Inches(2.85)
        stat_card(s, kx, ky, kw, Inches(1.0), lab, val, accent=col)


def m_portal(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "External Auditor Portal",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Read-only, scoped, expiring access for Big-4 teams",
             size=9, color=SLATE_500)
    rect(s, x, y + Inches(0.7), w - Inches(0.2), Inches(2.9),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, x + Inches(0.2), y + Inches(0.85), Inches(3),
             Inches(0.3), "PWC — SOX FY26 Walkthrough", size=11,
             bold=True, color=WHITE)
    severity_pill(s, x + Inches(4.6), y + Inches(0.85), "approved")
    text_box(s, x + Inches(0.2), y + Inches(1.2), Inches(5),
             Inches(0.25), "5 reviewers · expires 2026-09-30",
             size=9, color=SLATE_500)
    # access summary
    items = [("Engagements visible", "3"), ("Workpapers shared", "127"),
             ("Findings shared", "48"), ("Documents shared", "92"),
             ("Comments posted", "34"), ("Last login", "2h ago")]
    for i, (lab, val) in enumerate(items):
        r, c = divmod(i, 2)
        ix = x + Inches(0.2) + c * Inches(2.95)
        iy = y + Inches(1.65) + r * Inches(0.55)
        text_box(s, ix, iy, Inches(2), Inches(0.25), lab.upper(),
                 size=8, color=SLATE_500, bold=True)
        text_box(s, ix, iy + Inches(0.22), Inches(2), Inches(0.3), val,
                 size=14, bold=True, color=WHITE)


def m_charter(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Internal Audit Charter",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Versioned · AI-drafted · Committee-approved",
             size=9, color=SLATE_500)
    # Document mock
    rect(s, x, y + Inches(0.7), Inches(3.6), Inches(2.9),
         fill=WHITE, border=BORDER, corner=0.04)
    text_box(s, x + Inches(0.2), y + Inches(0.85), Inches(3.2),
             Inches(0.3), "INTERNAL AUDIT CHARTER",
             size=9, bold=True, color=RGBColor(0x33, 0x41, 0x55))
    text_box(s, x + Inches(0.2), y + Inches(1.15), Inches(3.2),
             Inches(0.3), "v3.2  ·  Adopted 14 Mar 2026",
             size=8, color=RGBColor(0x64, 0x74, 0x8B))
    for i in range(8):
        ly = y + Inches(1.55 + i * 0.22)
        line_w = Inches(3.0 if i % 3 != 2 else 2.4)
        rect(s, x + Inches(0.2), ly, line_w, Inches(0.07),
             fill=RGBColor(0xE2, 0xE8, 0xF0),
             border=RGBColor(0xE2, 0xE8, 0xF0), corner=0.5)
    # Right side: version timeline
    text_box(s, x + Inches(3.85), y + Inches(0.7), Inches(2.4),
             Inches(0.3), "VERSION HISTORY", size=9, bold=True,
             color=SLATE_500)
    versions = [("v3.2", "Approved", "Mar 2026", EMERALD),
                ("v3.1", "Superseded", "Sep 2025", SLATE_500),
                ("v3.0", "Superseded", "Mar 2024", SLATE_500),
                ("v2.4", "Archived", "Sep 2023", SLATE_500)]
    for i, (ver, st, dt, col) in enumerate(versions):
        vy = y + Inches(1.1 + i * 0.55)
        rect(s, x + Inches(3.85), vy, Inches(2.3), Inches(0.45),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, x + Inches(4.0), vy + Inches(0.08), Inches(0.6),
                 Inches(0.25), ver, size=10, bold=True, color=WHITE)
        text_box(s, x + Inches(4.0), vy + Inches(0.27), Inches(1.4),
                 Inches(0.2), st, size=8, color=col, bold=True)
        text_box(s, x + Inches(5.4), vy + Inches(0.13), Inches(0.8),
                 Inches(0.25), dt, size=8, color=SLATE_400,
                 align=PP_ALIGN.RIGHT)


def m_ccm(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3),
             "Continuous Controls Monitoring",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Daily automated tests · 28 controls · 6 frameworks",
             size=9, color=SLATE_500)
    # Daily strip
    days = 28
    sx = x; sy = y + Inches(0.75)
    cell = (w - Inches(0.4)) / days
    text_box(s, sx, sy, w, Inches(0.3), "MFA COVERAGE — LAST 28 DAYS",
             size=9, bold=True, color=SLATE_500)
    for i in range(days):
        v = (i * 7 + 3) % 28
        col = EMERALD if v < 22 else (AMBER if v < 26 else RED)
        cy = sy + Inches(0.35)
        rect(s, sx + Emu(int(cell * i)) + Inches(0.05),
             cy, cell - Inches(0.04), Inches(0.5),
             fill=col, border=col, corner=0.2)
    # Controls list
    text_box(s, x, y + Inches(1.65), w, Inches(0.3),
             "ACTIVE CONTROLS", size=9, bold=True, color=SLATE_500)
    rows = [
        ("CTL-014", "Privileged session reviews", "Pass", EMERALD, "100%"),
        ("CTL-022", "MFA enforcement (Cards)", "Fail", RED, "82%"),
        ("CTL-031", "Backup restore drill", "Warn", AMBER, "—"),
        ("CTL-047", "Vendor SOC2 currency", "Pass", EMERALD, "94%"),
    ]
    for i, (cid, name, st, col, val) in enumerate(rows):
        ry = y + Inches(2.0 + i * 0.42)
        rect(s, x, ry, w - Inches(0.2), Inches(0.36),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, x + Inches(0.15), ry + Inches(0.08), Inches(0.8),
                 Inches(0.25), cid, size=9, bold=True, color=SLATE_400)
        text_box(s, x + Inches(1.05), ry + Inches(0.08), Inches(2.6),
                 Inches(0.25), name, size=9, color=WHITE)
        pill(s, x + Inches(3.85), ry + Inches(0.06), st, fill=col, color=WHITE,
             w=Inches(0.7), h=Inches(0.24), size=8)
        text_box(s, x + Inches(4.7), ry + Inches(0.08), Inches(1.2),
                 Inches(0.25), val, size=9, bold=True, color=col,
                 align=PP_ALIGN.RIGHT)


def m_reporting(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Reporting — Audit Committee Pack",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "One-click generation · branded PDF · drill-through",
             size=9, color=SLATE_500)
    # Mock report cover
    rect(s, x, y + Inches(0.75), Inches(3.4), Inches(2.85),
         fill=PANEL, border=BLUE, border_w=2, corner=0.04)
    text_box(s, x + Inches(0.25), y + Inches(0.95), Inches(3),
             Inches(0.3), "AUDIT COMMITTEE", size=9, bold=True,
             color=BLUE)
    text_box(s, x + Inches(0.25), y + Inches(1.25), Inches(3),
             Inches(0.6), "Q1 FY26 Pack", size=18, bold=True, color=WHITE)
    text_box(s, x + Inches(0.25), y + Inches(1.85), Inches(3),
             Inches(0.3), "12 May 2026", size=9, color=SLATE_400)
    divider(s, x + Inches(0.25), y + Inches(2.25), Inches(2.9),
            color=BORDER_SOFT)
    sections = ["Executive summary", "Universe coverage",
                "Top findings", "Issue aging", "Plan vs. actual"]
    for i, sec in enumerate(sections):
        text_box(s, x + Inches(0.25), y + Inches(2.4 + i * 0.22),
                 Inches(3), Inches(0.25),
                 f"{i+1:02d}  {sec}", size=9, color=SLATE_300)
    # Right side: includes
    text_box(s, x + Inches(3.7), y + Inches(0.75), Inches(3),
             Inches(0.3), "AUTO-INCLUDED", size=9, bold=True,
             color=SLATE_500)
    chips_data = [("Heatmap snapshot", BLUE),
                  ("Top 10 findings (CCCE)", RED),
                  ("Aging buckets", AMBER),
                  ("Coverage delta vs. plan", EMERALD),
                  ("CCM exception summary", PURPLE),
                  ("Charter version reference", CYAN)]
    for i, (lab, col) in enumerate(chips_data):
        cy = y + Inches(1.1 + i * 0.42)
        rect(s, x + Inches(3.7), cy, Inches(2.5), Inches(0.34),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               x + Inches(3.85), cy + Inches(0.11),
                               Inches(0.13), Inches(0.13))
        d.fill.solid(); d.fill.fore_color.rgb = col
        d.line.fill.background()
        text_box(s, x + Inches(4.1), cy + Inches(0.07), Inches(2.2),
                 Inches(0.25), lab, size=9, color=WHITE)


def m_qaip(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3),
             "Quality Assurance & Improvement Program",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "IIA-aligned internal/external assessments + KPIs",
             size=9, color=SLATE_500)
    # KPI bar chart
    rect(s, x, y + Inches(0.7), w - Inches(0.2), Inches(1.9),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, x + Inches(0.2), y + Inches(0.8), Inches(3),
             Inches(0.3), "KPIs vs. target", size=9, bold=True,
             color=SLATE_400)
    kpis = [("Stakeholder satisfaction", 0.92, 0.85),
            ("Plan completion", 0.78, 0.85),
            ("Recommendations accepted", 0.88, 0.80),
            ("Cycle time", 0.72, 0.75)]
    base_y = y + Inches(1.15)
    bar_h = Inches(0.18)
    for i, (lab, val, tgt) in enumerate(kpis):
        ly = base_y + i * Inches(0.32)
        text_box(s, x + Inches(0.2), ly - Inches(0.02), Inches(2.8),
                 Inches(0.22), lab, size=8, color=SLATE_300)
        track_x = x + Inches(3.1)
        track_w = w - Inches(3.6)
        rect(s, track_x, ly + Inches(0.04), track_w, bar_h,
             fill=BG, border=BORDER_SOFT, corner=0.5)
        col = EMERALD if val >= tgt else AMBER
        rect(s, track_x, ly + Inches(0.04), Emu(int(track_w * val)),
             bar_h, fill=col, border=col, corner=0.5)
        # target marker
        tx = track_x + Emu(int(track_w * tgt))
        rect(s, tx, ly, Inches(0.03), Inches(0.27),
             fill=WHITE, border=WHITE)
        text_box(s, x + w - Inches(0.85), ly, Inches(0.6), Inches(0.25),
                 f"{int(val*100)}%", size=9, bold=True, color=col,
                 align=PP_ALIGN.RIGHT)
    # Last EQA
    sy = y + Inches(2.8)
    rect(s, x, sy, w - Inches(0.2), Inches(0.85),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, x + Inches(0.2), sy + Inches(0.12), Inches(3),
             Inches(0.3), "LAST EXTERNAL ASSESSMENT", size=9,
             bold=True, color=SLATE_500)
    text_box(s, x + Inches(0.2), sy + Inches(0.4), Inches(5),
             Inches(0.3), "IIA EQA — General Conformance",
             size=12, bold=True, color=WHITE)
    pill(s, x + w - Inches(2.0), sy + Inches(0.25),
         "Conformance", fill=EMERALD_SOFT, color=EMERALD,
         w=Inches(1.6), h=Inches(0.32), size=10)


def m_test_scripts(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3),
             "Test Script Library", size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "184 reusable scripts · auto-cloned into engagements",
             size=9, color=SLATE_500)
    # filter chips
    cy = y + Inches(0.7)
    chips_data = [("All (184)", BLUE), ("ITGC (54)", PURPLE),
                  ("AML (21)", AMBER), ("Financial (38)", EMERALD),
                  ("Cyber (47)", RED), ("Vendor (24)", CYAN)]
    cx = x
    for lab, col in chips_data:
        pill(s, cx, cy, lab, fill=PANEL_SOFT, color=col,
             w=Inches(1.05), h=Inches(0.32), size=9)
        cx += Inches(1.1)
    # script list
    rows = [
        ("TS-014", "Privileged access review walkthrough",
         "ITGC", "ISO 27001 A.5.18", PURPLE),
        ("TS-031", "MFA enforcement test of operating effectiveness",
         "Cyber", "NIST CSF PR.AC-7", RED),
        ("TS-058", "Vendor SOC2 currency check",
         "Vendor", "Internal VRM-04", CYAN),
        ("TS-072", "AML name-screening false-positive sample",
         "AML", "FATF R10", AMBER),
        ("TS-091", "Quarterly access recertification",
         "ITGC", "SOX ITGC", PURPLE),
    ]
    for i, (tid, name, cat, ref, col) in enumerate(rows):
        ry = y + Inches(1.25 + i * 0.45)
        rect(s, x, ry, w - Inches(0.2), Inches(0.4),
             fill=PANEL_SOFT, border=BORDER_SOFT)
        text_box(s, x + Inches(0.15), ry + Inches(0.1), Inches(0.7),
                 Inches(0.25), tid, size=9, bold=True, color=SLATE_400)
        text_box(s, x + Inches(0.95), ry + Inches(0.1), Inches(3.3),
                 Inches(0.25), name, size=9, color=WHITE)
        pill(s, x + Inches(4.3), ry + Inches(0.07), cat,
             fill=PANEL, color=col, w=Inches(0.7), h=Inches(0.26), size=8)
        text_box(s, x + Inches(5.1), ry + Inches(0.1), Inches(1.0),
                 Inches(0.25), ref, size=8, color=SLATE_500,
                 align=PP_ALIGN.RIGHT)


def m_skill_matrix(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Skill Matrix — Team Coverage",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Coverage by competency × seniority", size=9, color=SLATE_500)
    # Grid: 5 people x 6 skills heatmap
    people = ["P. Singh", "L. Chen", "A. Khan", "R. Mendoza", "T. Olsson"]
    skills = ["ITGC", "Data An.", "AML", "Cyber", "Vendor", "Audit Mgmt"]
    levels = [
        [4, 3, 2, 4, 2, 4],
        [3, 4, 1, 3, 3, 3],
        [2, 4, 4, 2, 1, 2],
        [4, 2, 3, 1, 4, 3],
        [3, 3, 4, 4, 2, 4],
    ]
    grid_x = x + Inches(1.1); grid_y = y + Inches(1.0)
    cell_w = (w - Inches(1.4)) / len(skills)
    cell_h = Inches(0.45)
    # column headers
    for c, sk in enumerate(skills):
        text_box(s, grid_x + c * cell_w, grid_y - Inches(0.3),
                 cell_w, Inches(0.25), sk, size=8, color=SLATE_400,
                 bold=True, align=PP_ALIGN.CENTER)
    for r, person in enumerate(people):
        text_box(s, x, grid_y + r * cell_h + Inches(0.1),
                 Inches(1.0), Inches(0.3), person, size=9, color=SLATE_300)
        for c, lvl in enumerate(levels[r]):
            cx = grid_x + c * cell_w
            cy = grid_y + r * cell_h
            level_color = {
                4: EMERALD, 3: BLUE, 2: AMBER, 1: RED
            }.get(lvl, SLATE_500)
            rect(s, cx + Inches(0.04), cy + Inches(0.04),
                 cell_w - Inches(0.08), cell_h - Inches(0.08),
                 fill=level_color, border=level_color, corner=0.18)
            text_box(s, cx, cy, cell_w, cell_h, str(lvl),
                     size=11, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # legend
    ly = y + Inches(3.5)
    leg = [(4, "Expert", EMERALD), (3, "Proficient", BLUE),
           (2, "Working", AMBER), (1, "Aware", RED)]
    for i, (n, lab, col) in enumerate(leg):
        lx = x + Inches(i * 1.5)
        rect(s, lx, ly, Inches(0.25), Inches(0.25), fill=col, border=col,
             corner=0.25)
        text_box(s, lx + Inches(0.3), ly, Inches(1.2), Inches(0.25),
                 f"{n} · {lab}", size=8, color=SLATE_400)


def m_capacity(s, x, y, w, h):
    text_box(s, x, y, w, Inches(0.3), "Capacity Planning",
             size=11, bold=True, color=WHITE)
    text_box(s, x, y + Inches(0.3), w, Inches(0.25),
             "Hours allocated · available · utilisation",
             size=9, color=SLATE_500)
    # Auditor utilisation bars
    rows = [
        ("P. Singh", 1640, 1820, BLUE),
        ("L. Chen", 1980, 1820, RED),
        ("A. Khan", 1410, 1820, EMERALD),
        ("R. Mendoza", 1750, 1820, AMBER),
        ("T. Olsson", 1530, 1820, EMERALD),
    ]
    for i, (name, used, avail, col) in enumerate(rows):
        ry = y + Inches(0.85 + i * 0.45)
        text_box(s, x, ry + Inches(0.05), Inches(1.4), Inches(0.3),
                 name, size=10, color=WHITE, bold=True)
        track_x = x + Inches(1.5)
        track_w = w - Inches(2.4)
        rect(s, track_x, ry + Inches(0.1), track_w, Inches(0.22),
             fill=BG, border=BORDER_SOFT, corner=0.5)
        pct = min(1.2, used / avail)
        fill_col = col if used <= avail else RED
        rect(s, track_x, ry + Inches(0.1),
             Emu(int(track_w * min(1.0, pct))), Inches(0.22),
             fill=fill_col, border=fill_col, corner=0.5)
        # marker at 100%
        text_box(s, x + w - Inches(0.85), ry + Inches(0.05),
                 Inches(0.7), Inches(0.3),
                 f"{int(pct*100)}%", size=10, bold=True, color=fill_col,
                 align=PP_ALIGN.RIGHT)
    # Summary chip
    sy = y + Inches(3.2)
    rect(s, x, sy, w - Inches(0.2), Inches(0.5),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, x + Inches(0.2), sy + Inches(0.13), w, Inches(0.25),
             "Forecast: 2 auditors over-allocated in Q3 — replan or hire",
             size=10, color=AMBER, bold=True)


def make_module_slides(prs):
    modules = [
        ("Universe", "Audit Universe",
         {"lead": "Risk-rated map of every auditable entity, refreshed continuously.",
          "body": "Auditable entities live in spreadsheets that are out of date the moment they're shared. Coverage gaps are invisible until the audit committee asks."},
         "Universe gives executives a live heat map across business units, scored for inherent risk and last-audited recency. Assurance gaps surface automatically.",
         m_universe,
         ["Tag entities by BU, geography, regulation",
          "Inherent + residual risk scoring",
          "Coverage delta vs. last 12/24/36 months",
          "Drill from heatmap to engagement"]),

        ("Plans", "Audit Plans",
         {"lead": "Risk-based annual planning, not a calendar exercise.",
          "body": "Annual plans get built in Excel, then drift the moment Q1 starts. There's no traceability from a plan slot back to the risk that justified it."},
         "Plans turns the universe into a multi-year, risk-weighted schedule with budgeted hours, dependencies, and live status for every engagement slot.",
         m_plans,
         ["Multi-year and rolling-12 views",
          "Auto-suggest from universe risk score",
          "Hours budgeted vs. capacity available",
          "Approval workflow with audit committee"]),

        ("Engagements", "Engagements",
         {"lead": "Run every engagement to one consistent shape.",
          "body": "Each lead auditor structures engagements differently. Reviewers can't tell, slot-to-slot, where time is being burned or where scope quietly grew."},
         "Engagements gives every audit a standard scope, objectives, milestones, hours and team — with AI scope-drafting from the universe entry.",
         m_engagements,
         ["AI-drafted scope from entity context",
          "Budget vs. actual hours, live",
          "Standard milestones + workpaper structure",
          "Assurance gap badges from universe"]),

        ("Workpapers", "Workpapers & Sign-Off",
         {"lead": "Test procedures, evidence and review trail in one place.",
          "body": "Workpapers live in Word, evidence in shared drives, review notes in email. Reconstructing what happened months later is forensic work."},
         "Workpapers structures every test (planning, walkthrough, analytical, interview) with a 4-stage sign-off chain: prepare → review → lead sign-off → close.",
         m_workpapers,
         ["Standard procedure templates",
          "Inline evidence + cross-reference",
          "Preparer / reviewer / lead signatures",
          "Reject-and-return with notes"]),

        ("Findings", "Findings",
         {"lead": "Every finding written to the same CCCE structure.",
          "body": "Findings get rewritten by every reviewer because there's no shared template. Severity ratings drift between auditors. Recurring issues go undetected."},
         "Findings enforces the IIA-aligned CCCE structure (Condition, Criteria, Cause, Effect), with AI-drafted condition statements and severity calibration.",
         m_findings,
         ["AI-drafted CCCE from anomaly description",
          "Severity calibration across the team",
          "Similar / recurring detection across years",
          "Bulk import + Excel template"]),

        ("Issues", "Issue Tracking & Escalation",
         {"lead": "No more 'we'll chase that one next month'.",
          "body": "Once a finding has a due date, ownership tracking ends. Overdue issues sit unresolved because no one knows whose desk they're on."},
         "Issues runs aging buckets, automatic 3-level escalation (reminder → manager → executive), and a single dashboard for everything overdue.",
         m_issues,
         ["Aging buckets: 0-30, 30-60, 60-90, 90+ days",
          "Auto-escalate from L1 → L2 → L3",
          "Daily digest to action owners",
          "Resolved-with-notes audit trail"]),

        ("Surveys", "Surveys & Self-Assessments",
         {"lead": "Pre-audit, post-audit, and CSA without SurveyMonkey.",
          "body": "Pre-audit questionnaires get sent in Outlook with a Word attachment. Responses arrive scattered, half complete, and need to be hand-tallied."},
         "Surveys lets the team build pre-audit, post-audit, and control self-assessment surveys — with AI question generation tied to the engagement scope.",
         m_surveys,
         ["AI-generated questions from scope",
          "Pre/post-audit + CSA templates",
          "Email distribution + reminders",
          "Response dashboard + export"]),

        ("Documents", "Document Repository",
         {"lead": "One place for evidence, working papers, reports, charters.",
          "body": "Audit evidence lives in OneDrive, SharePoint, email and personal laptops. Confidentiality is enforced by hope. Retention is whatever default the drive has."},
         "Documents centralises every artefact with type, engagement linkage, confidentiality flag and full-text search — backed by the same access model as the rest of the platform.",
         m_documents,
         ["Tag by type, engagement, BU",
          "Confidentiality flag with access scope",
          "Full-text search across the corpus",
          "Linked to workpapers + findings"]),

        ("Analytics", "Audit Analytics",
         {"lead": "The numbers behind the assurance — not a quarterly slide.",
          "body": "Analytics today means rebuilding the same Excel pivot at the end of every quarter. Trends are guessed at; hard numbers arrive too late to act on."},
         "Analytics ships a live coverage / velocity / closure / aging dashboard, sliced by BU, engagement type and severity — with drill-through to underlying records.",
         m_analytics,
         ["Coverage % by BU and risk tier",
          "Mean cycle time + closure rate",
          "Findings opened vs. closed by quarter",
          "Slice by severity, owner, theme"]),

        ("Portal", "External Auditor Portal",
         {"lead": "Give Big-4 a workspace, not a Dropbox link.",
          "body": "External auditors get a shared folder and a user account that's hard to revoke. Read-only is enforced by access groups that nobody audits."},
         "Portal gives external teams scoped, expiring, read-only access to exactly the engagements / workpapers / findings they need — with full activity logging.",
         m_portal,
         ["Scoped to specific engagements",
          "Expiry date + auto-revoke",
          "Read-only with comment threads",
          "Full activity log per external user"]),

        ("Charter", "Audit Charter",
         {"lead": "Living charter, versioned, approved.",
          "body": "The audit charter lives as a Word doc on the head of audit's laptop. Nobody knows which version is current. Committee approval is hunted in old minutes."},
         "Charter holds the current and historical charter documents, with AI-assisted drafting, version comparison, and committee approval workflow.",
         m_charter,
         ["AI drafting from IIA template",
          "Side-by-side version diff",
          "Committee approval workflow",
          "Linked from every engagement"]),

        ("CCM", "Continuous Controls Monitoring",
         {"lead": "Daily controls testing, not annual snapshots.",
          "body": "Most controls get tested once a year on a sample. Failures between tests go undetected for months. Internal audit and SecOps don't share data."},
         "CCM runs scripted controls daily against source systems, surfaces exceptions, and feeds them into the audit findings backlog automatically.",
         m_ccm,
         ["Scripted controls per framework",
          "Daily execution + alerting",
          "Exception → finding workflow",
          "Time-series view per control"]),

        ("Reporting", "Reporting",
         {"lead": "The committee pack writes itself.",
          "body": "Quarterly reporting is two analysts and a week of PowerPoint. By the time it's done, half the numbers are stale and the layout has drifted from last quarter."},
         "Reporting generates the committee pack from live data — heatmap, top findings, aging, coverage delta — in a branded template, in one click.",
         m_reporting,
         ["Branded PDF, one-click export",
          "Auto-refreshed snapshots",
          "Drill-through to workpapers",
          "Historical pack archive"]),

        ("QAIP", "Quality Assurance & Improvement",
         {"lead": "IIA-aligned QAIP that runs itself.",
          "body": "QAIP is the program everyone agrees they need and nobody runs. Internal assessments are sporadic; external EQAs surface preventable findings."},
         "QAIP holds internal/external assessment cycles, KPI tracking against IIA standards, and the conformance evidence the committee needs.",
         m_qaip,
         ["Internal assessment scheduling",
          "KPIs vs. target with targets visible",
          "External EQA history + actions",
          "Conformance statement export"]),

        ("Test Scripts", "Test Script Library",
         {"lead": "184 reusable scripts. Cloned, not retyped.",
          "body": "Every engagement re-writes the same MFA test, vendor SOC2 check, and access recertification procedure from memory. Quality varies wildly."},
         "Test Scripts is a versioned library of reusable procedures, mapped to frameworks and categories, that auditors clone into a workpaper in two clicks.",
         m_test_scripts,
         ["Versioned, peer-reviewed library",
          "Mapped to ISO, NIST, SOX, FATF",
          "Clone-to-workpaper in two clicks",
          "Category + framework filters"]),

        ("Skill Matrix", "Skill Matrix",
         {"lead": "Match the right auditor to the right engagement.",
          "body": "Engagement staffing is done from memory. Junior auditors get put on cyber engagements they aren't ready for. Skills gaps surface only when work slips."},
         "Skill Matrix tracks every auditor's competency × seniority, surfaces team-wide gaps, and suggests staffing for incoming engagements.",
         m_skill_matrix,
         ["Competency × seniority heatmap",
          "Gap detection at team level",
          "Staffing suggestions per engagement",
          "Training-plan input"]),

        ("Capacity", "Capacity Planning",
         {"lead": "See over-allocation before it ships.",
          "body": "Audit teams routinely commit to more hours than they have. The over-allocation surfaces in Q3 when something has to be deferred or descoped."},
         "Capacity reconciles planned engagement hours against auditor availability — daily, monthly, quarterly — and flags over-allocation before it bites.",
         m_capacity,
         ["Hours planned vs. available",
          "Per-auditor utilisation",
          "Forecast 1-2 quarters ahead",
          "Replan or hire signal"]),
    ]
    return modules


def make_roi(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    text_box(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
             "BUSINESS CASE", size=10, bold=True, color=EMERALD)
    text_box(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9),
             "What AuditVerse.AI means in numbers.",
             font=FONT_DISPLAY, size=30, bold=True, color=WHITE)
    text_box(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             "Indicative impact for a 12-person internal audit function.",
             size=12, color=SLATE_400)

    big_stats = [
        ("40%", "less cycle time", "per engagement, planning to report",
         EMERALD),
        ("60%", "less admin", "writing/formatting, not auditing", BLUE),
        ("3x", "audit-committee freshness",
         "live numbers vs. quarterly pack", PURPLE),
        ("100%", "evidence retrievability",
         "every workpaper indexed and signed off", AMBER),
    ]
    for i, (val, lab, sub, col) in enumerate(big_stats):
        x = Inches(0.6 + i * 3.05)
        rect(s, x, Inches(2.3), Inches(2.95), Inches(2.4),
             fill=PANEL, border=BORDER)
        text_box(s, x + Inches(0.2), Inches(2.5), Inches(2.6),
                 Inches(1.3), val, font=FONT_DISPLAY, size=58, bold=True,
                 color=col, line_spacing=0.95)
        text_box(s, x + Inches(0.2), Inches(3.85), Inches(2.6),
                 Inches(0.4), lab, size=14, bold=True, color=WHITE)
        text_box(s, x + Inches(0.2), Inches(4.2), Inches(2.6),
                 Inches(0.5), sub, size=10, color=SLATE_400,
                 line_spacing=1.3)

    # Bottom: assumptions row
    rect(s, Inches(0.6), Inches(5.0), Inches(12.15), Inches(1.7),
         fill=PANEL_SOFT, border=BORDER_SOFT)
    text_box(s, Inches(0.85), Inches(5.15), Inches(8), Inches(0.3),
             "WHAT MAKES THE NUMBERS HOLD", size=10, bold=True,
             color=SLATE_500)
    points = [
        ("AI scope drafting", "2-4 hrs / engagement back to the auditor"),
        ("CCCE structuring", "30 min / finding back to the lead"),
        ("Auto-escalation", "Issues never sit > 30 days in silence"),
        ("Live committee pack", "0 prep weeks per quarter"),
    ]
    for i, (h, s_text) in enumerate(points):
        px = Inches(0.85 + i * 2.95)
        text_box(s, px, Inches(5.55), Inches(2.8), Inches(0.3),
                 h, size=11, bold=True, color=BLUE)
        text_box(s, px, Inches(5.85), Inches(2.8), Inches(0.7),
                 s_text, size=10, color=SLATE_300, line_spacing=1.35)
    footer(s, total - 1, total)


def make_cta(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    # Two-column layout: left CTA, right contact card
    text_box(s, Inches(0.7), Inches(1.0), Inches(2), Inches(0.35),
             "NEXT STEPS", size=10, bold=True, color=BLUE)
    text_box(s, Inches(0.7), Inches(1.4), Inches(7.5), Inches(2.4),
             "Let's run\nyour next audit\non AuditVerse.AI.",
             font=FONT_DISPLAY, size=52, bold=True, color=WHITE,
             line_spacing=0.95)
    text_box(s, Inches(0.7), Inches(4.7), Inches(7), Inches(1),
             ("90-minute pilot scoping. We bring 2 reference engagements; "
              "you bring your real universe and a recent finding."),
             size=14, color=SLATE_300, line_spacing=1.4)
    # Three-step strip
    steps = [
        ("01", "Scoping call", "60 min · your team + ours"),
        ("02", "Pilot", "30 days · 2 engagements"),
        ("03", "Decision", "Roll out or walk away"),
    ]
    for i, (n, h, sub) in enumerate(steps):
        sx = Inches(0.7 + i * 2.3)
        rect(s, sx, Inches(5.85), Inches(2.15), Inches(1.0),
             fill=PANEL, border=BORDER)
        text_box(s, sx + Inches(0.2), Inches(5.95), Inches(0.5),
                 Inches(0.3), n, size=11, bold=True, color=BLUE)
        text_box(s, sx + Inches(0.7), Inches(5.95), Inches(1.4),
                 Inches(0.3), h, size=11, bold=True, color=WHITE)
        text_box(s, sx + Inches(0.7), Inches(6.25), Inches(1.4),
                 Inches(0.5), sub, size=8, color=SLATE_400,
                 line_spacing=1.3)

    # Right contact panel
    rect(s, Inches(8.5), Inches(1.0), Inches(4.3), Inches(5.85),
         fill=PANEL, border=BORDER)
    text_box(s, Inches(8.8), Inches(1.3), Inches(3.7), Inches(0.35),
             "AUDITVERSE.AI", size=10, bold=True, color=BLUE)
    text_box(s, Inches(8.8), Inches(1.65), Inches(3.7), Inches(0.6),
             "Audit Management Platform", size=18, bold=True, color=WHITE)
    divider(s, Inches(8.8), Inches(2.3), Inches(3.7), color=BORDER_SOFT)
    blocks = [
        ("PILOT CONTACT", "pilot@auditverse.ai"),
        ("SOLUTIONS", "solutions@auditverse.ai"),
        ("WEB", "auditverse.ai"),
        ("PHONE", "+1 (555) 014-0210"),
    ]
    for i, (lab, val) in enumerate(blocks):
        by = Inches(2.6 + i * 0.95)
        text_box(s, Inches(8.8), by, Inches(3.7), Inches(0.3),
                 lab, size=9, bold=True, color=SLATE_500)
        text_box(s, Inches(8.8), by + Inches(0.3), Inches(3.7), Inches(0.4),
                 val, size=14, bold=True, color=WHITE)
    footer(s, total, total)


# ---------- Build ----------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    modules = make_module_slides(prs)
    total_slides = 1 + 1 + 1 + len(modules) + 1 + 1  # cover + problem + solution + modules + roi + cta

    make_cover(prs)
    make_problem(prs)
    make_solution(prs)
    for i, mod in enumerate(modules):
        eyebrow, title, problem, solution, mockup_fn, features = mod
        module_slide(prs, eyebrow, title, problem, solution, mockup_fn,
                     features, page_num=i + 4, total=total_slides)
    make_roi(prs, total_slides)
    make_cta(prs, total_slides)

    out = Path("assets/pitch")
    out.mkdir(parents=True, exist_ok=True)
    pptx_path = out / "AuditVerse.AI_Pitch_Deck.pptx"
    prs.save(pptx_path)
    print(f"Wrote {pptx_path}  ({total_slides} slides)")
    return pptx_path


if __name__ == "__main__":
    build()
